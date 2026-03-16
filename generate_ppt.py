"""Generate NORA EOS presentation PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import os
ASSETS = r"c:\Users\Shaggy\echo\ppt_assets"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_image(slide, left, top, width, height, filename):
    path = os.path.join(ASSETS, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        placeholder_box(slide, left, top, width, height, f"[Missing: {filename}]")

# ── Color palette ──
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT = RGBColor(0xE5, 0xE7, 0xEB)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT, bullet_color=ACCENT_CYAN, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_CYAN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def placeholder_box(slide, left, top, width, height, label="[Screenshot]"):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
    shape.line.color.rgb = RGBColor(0x37, 0x41, 0x51)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(4.5), Inches(3.1), Inches(4.3), ACCENT_CYAN)
add_text(slide, Inches(0), Inches(1.8), Inches(13.333), Inches(1.2),
         "NORA EOS", 54, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.8),
         "Enterprise Operating System", 28, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(4.3), Inches(13.333), Inches(0.6),
         "ECHO  \u2022  HERA  \u2022  ARGUS", 20, ACCENT_CYAN, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS NORA EOS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "What is NORA EOS?", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_CYAN)

add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(2), [
    "An AI-powered enterprise operating system built for modern teams",
    "Three integrated modules that automate communication, project management, and performance tracking",
    "Connects Gmail, Google Calendar, and internal task management into one intelligent layer",
    "Every module feeds data to the others \u2014 creating a closed-loop productivity system",
], 18, LIGHT)

# Three module cards
modules = [
    ("ECHO", "Smart Email Assistant", ACCENT_CYAN, "AI email classification, reply suggestions, calendar management, and chat-driven drafting"),
    ("HERA", "Project Orchestrator", ACCENT_PURPLE, "Natural language project creation, Kanban board, Gantt timeline, and team capacity management"),
    ("ARGUS", "Productivity Intelligence", ACCENT_GREEN, "Automated scoring across 4 metrics, AI recommendations, and transparent performance tracking"),
]
for i, (name, subtitle, color, desc) in enumerate(modules):
    x = Inches(0.8 + i * 4.1)
    add_shape_bg(slide, x, Inches(4.0), Inches(3.7), Inches(3.0), BG_CARD, 0.05)
    add_accent_line(slide, x + Inches(0.2), Inches(4.2), Inches(1.5), color)
    add_text(slide, x + Inches(0.2), Inches(4.4), Inches(3.3), Inches(0.5),
             name, 24, color, True)
    add_text(slide, x + Inches(0.2), Inches(4.9), Inches(3.3), Inches(0.4),
             subtitle, 14, GRAY)
    add_text(slide, x + Inches(0.2), Inches(5.4), Inches(3.3), Inches(1.4),
             desc, 13, LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE OVERVIEW
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "System Architecture", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_CYAN)

add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), [
    "Each module runs as an independent microservice with its own database",
    "Cross-service communication via authenticated REST APIs",
    "Shared Celery + Redis worker pipelines for background processing",
    "LLM integration (Groq / Llama 3.3 70B) powers AI features across all modules",
    "Gmail and Google Calendar APIs provide real-time email and scheduling data",
], 16, LIGHT)

# Connection diagram placeholder
add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.5), BG_CARD, 0.03)

# ECHO box
add_shape_bg(slide, Inches(7.5), Inches(1.8), Inches(4.5), Inches(1.2), RGBColor(0x0A, 0x1E, 0x3D), 0.05)
add_text(slide, Inches(7.7), Inches(1.9), Inches(2), Inches(0.4), "ECHO", 18, ACCENT_CYAN, True)
add_text(slide, Inches(7.7), Inches(2.3), Inches(4), Inches(0.5), "Gmail API  \u2022  Calendar API  \u2022  Chat LLM", 11, GRAY)

# HERA box
add_shape_bg(slide, Inches(7.5), Inches(3.3), Inches(4.5), Inches(1.2), RGBColor(0x1A, 0x0E, 0x3D), 0.05)
add_text(slide, Inches(7.7), Inches(3.4), Inches(2), Inches(0.4), "HERA", 18, ACCENT_PURPLE, True)
add_text(slide, Inches(7.7), Inches(3.8), Inches(4), Inches(0.5), "Orchestration LLM  \u2022  Task DB  \u2022  Scheduling", 11, GRAY)

# ARGUS box
add_shape_bg(slide, Inches(7.5), Inches(4.8), Inches(4.5), Inches(1.2), RGBColor(0x0A, 0x2E, 0x1E), 0.05)
add_text(slide, Inches(7.7), Inches(4.9), Inches(2), Inches(0.4), "ARGUS", 18, ACCENT_GREEN, True)
add_text(slide, Inches(7.7), Inches(5.3), Inches(4), Inches(0.5), "Scoring Engine  \u2022  Advice LLM  \u2022  Metrics API", 11, GRAY)

# Arrows labels
add_text(slide, Inches(8.5), Inches(3.05), Inches(4), Inches(0.3),
         "\u2191\u2193  Notifications + Team Lookup", 10, GRAY)
add_text(slide, Inches(8.5), Inches(4.55), Inches(4), Inches(0.3),
         "\u2191\u2193  Metrics Collection", 10, GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — ECHO TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(4.5), Inches(3.1), Inches(4.3), ACCENT_CYAN)
add_text(slide, Inches(0), Inches(1.8), Inches(13.333), Inches(1.2),
         "ECHO", 54, ACCENT_CYAN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.8),
         "Smart Email Assistant", 28, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(4.5), Inches(13.333), Inches(0.6),
         "AI-powered email management, reply suggestions, and calendar intelligence", 16, LIGHT, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — ECHO CORE FEATURES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ECHO \u2014 Core Features", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_CYAN)

features = [
    ("Email Classification", "Incoming emails are automatically categorized as Work, Personal, Meeting, Newsletter, or Urgent using AI"),
    ("AI Reply Suggestions", "ECHO drafts personalized replies based on the user\u2019s writing style \u2014 users accept, edit, or reject"),
    ("Chat-Driven Email Drafting", "Users ask ECHO to draft emails via chat \u2014 recipients are resolved from contacts or HERA project teams"),
    ("Calendar Intelligence", "Direct event creation from chat, conflict detection, and automated meeting proposals with confirmation flow"),
    ("Cross-Module Notifications", "Task assignments from HERA appear as notifications \u2014 deadline reminders trigger automatically"),
    ("Daily Digest", "AI-generated summary of email activity, urgency breakdown, and suggestion acceptance stats"),
]
for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.6 + row * 1.8)
    add_shape_bg(slide, x, y, Inches(5.8), Inches(1.5), BG_CARD, 0.04)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(5.3), Inches(0.4),
             title, 17, ACCENT_CYAN, True)
    add_text(slide, x + Inches(0.25), y + Inches(0.6), Inches(5.3), Inches(0.8),
             desc, 13, LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — ECHO USER FLOW
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ECHO \u2014 User Flow", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_CYAN)

steps = [
    ("1", "Google Login", "OAuth authentication grants access to Gmail and Calendar"),
    ("2", "Onboarding", "6-step walkthrough \u2014 fetches and classifies 20 emails on completion"),
    ("3", "Dashboard", "Summary card, next meeting, email cards with AI suggestions, task cards"),
    ("4", "Chat", "Ask anything \u2014 ECHO responds with text, email drafts, or calendar events"),
    ("5", "Email Actions", "Accept suggestion, write custom reply, or draft from chat and send"),
    ("6", "Calendar", "Create events, detect conflicts, send meeting proposals with auto-confirmation"),
]
for i, (num, title, desc) in enumerate(steps):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.6 + row * 2.7)
    add_shape_bg(slide, x, y, Inches(3.7), Inches(2.2), BG_CARD, 0.04)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_CYAN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.85), y + Inches(0.25), Inches(2.6), Inches(0.4),
             title, 18, WHITE, True)
    add_text(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.3), Inches(1.2),
             desc, 13, LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — ECHO ARCHITECTURE DIAGRAM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ECHO \u2014 Architecture Diagram", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_CYAN)
add_image(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5), "arch_echo.png")

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — ECHO SCREENSHOTS (Login + Onboarding + Dashboard)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ECHO \u2014 Dashboard, Tasks & Daily Digest", 32, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.5), ACCENT_CYAN)

add_image(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(5.2), "echo_dashboard.png")
add_text(slide, Inches(0.5), Inches(6.85), Inches(3.8), Inches(0.4),
         "Dashboard with summary card and email stats", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(4.75), Inches(1.6), Inches(3.8), Inches(5.2), "echo_tasks.png")
add_text(slide, Inches(4.75), Inches(6.85), Inches(3.8), Inches(0.4),
         "Tasks tab with HERA-synced task cards", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(9), Inches(1.6), Inches(3.8), Inches(5.2), "echo_digest.png")
add_text(slide, Inches(9), Inches(6.85), Inches(3.8), Inches(0.4),
         "AI-generated daily email digest", 12, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — ECHO SCREENSHOTS (AI Suggestions + Chat + Calendar)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ECHO \u2014 Chat Email Drafting & Sending", 32, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.5), ACCENT_CYAN)

add_image(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.5), "echo_draft.png")
add_text(slide, Inches(0.5), Inches(7.05), Inches(3.8), Inches(0.4),
         "AI-drafted email with Send / Alter buttons", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(4.75), Inches(1.5), Inches(3.8), Inches(5.5), "echo_draft_team.png")
add_text(slide, Inches(4.75), Inches(7.05), Inches(3.8), Inches(0.4),
         "Draft to project team (resolved via HERA)", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(9), Inches(1.5), Inches(3.8), Inches(5.5), "echo_sent.png")
add_text(slide, Inches(9), Inches(7.05), Inches(3.8), Inches(0.4),
         "Email sent confirmation via Gmail API", 12, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — HERA TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(4.5), Inches(3.1), Inches(4.3), ACCENT_PURPLE)
add_text(slide, Inches(0), Inches(1.8), Inches(13.333), Inches(1.2),
         "HERA", 54, ACCENT_PURPLE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.8),
         "Project Orchestration Engine", 28, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(4.5), Inches(13.333), Inches(0.6),
         "Natural language project creation, real-time task tracking, and team management", 16, LIGHT, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — HERA CORE FEATURES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "HERA \u2014 Core Features", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_PURPLE)

features = [
    ("AI Orchestration", "Describe a project in plain English or upload documents \u2014 HERA generates tasks, assigns team members, and builds a sprint plan"),
    ("Kanban Board", "Real-time board with To Do, In Progress, and Done columns \u2014 updates every 5 seconds as employees work"),
    ("Gantt Timeline", "Visual schedule with task bars, dependency tracking, today marker, and overdue/blocked indicators"),
    ("Team Management", "Member cards with workload meters, skill badges, and capacity limits to prevent over-allocation"),
    ("Employee Portal", "Employees see only their assigned tasks \u2014 start, complete, and track deadlines independently"),
    ("Cross-Module Sync", "Task assignments trigger ECHO notifications \u2014 ECHO queries HERA for project team lookups"),
]
for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.6 + row * 1.8)
    add_shape_bg(slide, x, y, Inches(5.8), Inches(1.5), BG_CARD, 0.04)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(5.3), Inches(0.4),
             title, 17, ACCENT_PURPLE, True)
    add_text(slide, x + Inches(0.25), y + Inches(0.6), Inches(5.3), Inches(0.8),
             desc, 13, LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — HERA ADMIN FLOW
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "HERA \u2014 Admin Flow", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_PURPLE)

steps = [
    ("1", "Login", "Manager enters email and password"),
    ("2", "Orchestrate", "Describe project or upload docs \u2014 AI generates task breakdown"),
    ("3", "Review", "Adjust assignments, add/delete tasks, confirm with Start Project"),
    ("4", "Board", "Real-time Kanban \u2014 click any task to edit status, assignee, priority"),
    ("5", "Timeline", "Gantt chart with dependencies, deadlines, and blocked task indicators"),
    ("6", "Team", "Manage members, skills, and workload capacity"),
]
for i, (num, title, desc) in enumerate(steps):
    col = i % 6
    x = Inches(0.4 + col * 2.1)
    y = Inches(1.8)
    add_shape_bg(slide, x, y, Inches(1.9), Inches(2.8), BG_CARD, 0.04)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y + Inches(0.2), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_PURPLE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(slide, x + Inches(0.1), y + Inches(0.9), Inches(1.7), Inches(0.4),
             title, 16, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(1.4), Inches(1.7), Inches(1.2),
             desc, 12, LIGHT, False, PP_ALIGN.CENTER)

# Arrow connectors between steps
for i in range(5):
    x = Inches(0.4 + (i+1) * 2.1 - 0.15)
    add_text(slide, x, Inches(2.7), Inches(0.3), Inches(0.3), "\u2192", 20, ACCENT_PURPLE, True, PP_ALIGN.CENTER)

# Employee flow below
add_text(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.5),
         "Employee Flow", 24, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(5.45), Inches(1.8), ACCENT_PURPLE)

emp_steps = [
    ("Login", "Enter email only"),
    ("My Tasks", "View assigned tasks with priorities and deadlines"),
    ("Execute", "Start \u2192 In Progress \u2192 Complete"),
    ("Timeline", "Personal Gantt view of task schedule"),
]
for i, (title, desc) in enumerate(emp_steps):
    x = Inches(0.8 + i * 3.1)
    y = Inches(5.8)
    add_shape_bg(slide, x, y, Inches(2.7), Inches(1.2), BG_CARD, 0.04)
    add_text(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.4), Inches(0.35),
             title, 15, ACCENT_PURPLE, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.4), Inches(0.6),
             desc, 12, LIGHT, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 13 — HERA ARCHITECTURE DIAGRAM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "HERA \u2014 Architecture Diagram", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_PURPLE)
add_image(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5), "arch_hera.png")

# ════════════════════════════════════════════════════════════════
# SLIDE 14 — HERA SCREENSHOTS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "HERA \u2014 Orchestrate & Board", 32, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3), ACCENT_PURPLE)

add_image(slide, Inches(0.3), Inches(1.5), Inches(6.3), Inches(2.8), "hera_orchestrate.png")
add_text(slide, Inches(0.3), Inches(4.35), Inches(6.3), Inches(0.4),
         "AI project generation with team panel", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(6.9), Inches(1.5), Inches(6.3), Inches(2.8), "hera_board.png")
add_text(slide, Inches(6.9), Inches(4.35), Inches(6.3), Inches(0.4),
         "Real-time Kanban board with task cards", 12, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 15 — HERA SCREENSHOTS (Timeline + Team + Employee)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "HERA \u2014 Timeline, Team & Employee Portal", 32, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3.5), ACCENT_PURPLE)

add_image(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.2), "hera_timeline.png")
add_text(slide, Inches(0.5), Inches(6.85), Inches(3.8), Inches(0.4),
         "Gantt timeline with dependencies", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(4.75), Inches(1.5), Inches(3.8), Inches(5.2), "hera_team.png")
add_text(slide, Inches(4.75), Inches(6.85), Inches(3.8), Inches(0.4),
         "Member cards with workload and skills", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(9), Inches(1.5), Inches(3.8), Inches(5.2), "hera_mytasks.png")
add_text(slide, Inches(9), Inches(6.85), Inches(3.8), Inches(0.4),
         "Employee task view with Start/Done actions", 12, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 16 — ARGUS TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(4.5), Inches(3.1), Inches(4.3), ACCENT_GREEN)
add_text(slide, Inches(0), Inches(1.8), Inches(13.333), Inches(1.2),
         "ARGUS", 54, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.8),
         "Productivity Intelligence", 28, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(4.5), Inches(13.333), Inches(0.6),
         "Automated performance scoring, AI recommendations, and transparent metrics", 16, LIGHT, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 17 — ARGUS SCORING MODEL
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ARGUS \u2014 Scoring Model", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_GREEN)

# Formula
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.2), BG_CARD, 0.03)
add_text(slide, Inches(1.2), Inches(1.75), Inches(11), Inches(0.8),
         "Overall Score = Task Completion \u00d7 0.40 + Timeliness \u00d7 0.25 + Communication \u00d7 0.20 + Engagement \u00d7 0.15",
         18, ACCENT_GREEN, True, PP_ALIGN.CENTER)

metrics = [
    ("Task Completion", "40%", "Completion rate + priority-weighted output from HERA", ACCENT_GREEN),
    ("Timeliness", "25%", "On-time delivery rate for tasks with deadlines", ACCENT_CYAN),
    ("Communication", "20%", "Email suggestion acceptance rate from ECHO", ACCENT_PURPLE),
    ("Engagement", "15%", "Notification read and action rate from ECHO", RGBColor(0xF5, 0x9E, 0x0B)),
]
for i, (name, weight, desc, color) in enumerate(metrics):
    x = Inches(0.8 + i * 3.05)
    y = Inches(3.2)
    add_shape_bg(slide, x, y, Inches(2.8), Inches(2.5), BG_CARD, 0.04)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.4), Inches(0.4),
             name, 16, color, True, PP_ALIGN.CENTER)
    # Weight circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.7), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = weight
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, x + Inches(0.15), y + Inches(1.8), Inches(2.5), Inches(0.6),
             desc, 11, LIGHT, False, PP_ALIGN.CENTER)

# Data sources note
add_text(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
         "All data sourced transparently from HERA and ECHO \u2014 no subjective factors",
         14, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 18 — ARGUS FEATURES + FLOWS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ARGUS \u2014 Two User Paths", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_GREEN)

# Manager path
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.4), BG_CARD, 0.03)
add_text(slide, Inches(1.1), Inches(1.75), Inches(5), Inches(0.5),
         "Manager Dashboard", 22, ACCENT_GREEN, True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.2), Inches(4.5), [
    "Four stat cards: employee count, team average, top performer, needs support",
    "Category averages: visual bar charts across all 4 scoring dimensions",
    "Team leaderboard: color-coded scores with trend arrows (improving/declining/stable)",
    "Click any employee to drill into their detail page",
    "Employee detail: score gauge, breakdown bars, AI recommendations, score history",
    "AI advice cards with High / Medium / Low priority levels",
], 14, LIGHT, ACCENT_GREEN)

# Employee path
add_shape_bg(slide, Inches(7), Inches(1.6), Inches(5.8), Inches(5.4), BG_CARD, 0.03)
add_text(slide, Inches(7.3), Inches(1.75), Inches(5), Inches(0.5),
         "Employee \u2014 My Score", 22, ACCENT_GREEN, True)
add_bullet_list(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(4.5), [
    "Personal score gauge with color-coded performance label",
    "Score breakdown: 4 horizontal bars showing each metric weight",
    "\"How is my score calculated?\" \u2014 expandable formula with full transparency",
    "Data sources clearly labeled: which metric comes from HERA vs ECHO",
    "Dismissible AI recommendations personalized to their weak areas",
    "Score history table to track personal trends over time",
], 14, LIGHT, ACCENT_GREEN)

# ════════════════════════════════════════════════════════════════
# SLIDE 19 — ARGUS ARCHITECTURE DIAGRAM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ARGUS \u2014 Architecture Diagram", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_GREEN)
add_image(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5), "arch_argus.png")

# ════════════════════════════════════════════════════════════════
# SLIDE 20 — ARGUS SCREENSHOTS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ARGUS \u2014 Screenshots", 32, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_GREEN)

add_image(slide, Inches(0.5), Inches(1.5), Inches(6.1), Inches(5.5), "argus_dashboard.png")
add_text(slide, Inches(0.5), Inches(7.05), Inches(6.1), Inches(0.4),
         "Manager dashboard: team stats, leaderboard, AI recommendations", 12, GRAY, False, PP_ALIGN.CENTER)

add_image(slide, Inches(6.9), Inches(1.5), Inches(6.1), Inches(5.5), "argus_myscore.png")
add_text(slide, Inches(6.9), Inches(7.05), Inches(6.1), Inches(0.4),
         "Employee My Score: gauge, breakdown, and formula transparency", 12, GRAY, False, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 21 — CROSS-MODULE INTEGRATION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "Cross-Module Integration", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_CYAN)

connections = [
    ("HERA \u2192 ECHO", "Task assignments in HERA trigger real-time notifications in ECHO for assigned employees", ACCENT_PURPLE),
    ("ECHO \u2192 HERA", "ECHO queries HERA to resolve project team members when drafting emails or creating calendar events", ACCENT_CYAN),
    ("HERA \u2192 ARGUS", "Task completion rates and deadline adherence feed into ARGUS for productivity scoring", ACCENT_PURPLE),
    ("ECHO \u2192 ARGUS", "Suggestion acceptance rates and notification engagement feed into ARGUS scoring", ACCENT_CYAN),
    ("ARGUS \u2192 Managers", "AI-generated insights help managers make informed decisions about task assignments in HERA", ACCENT_GREEN),
]
for i, (label, desc, color) in enumerate(connections):
    y = Inches(1.6 + i * 1.1)
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.9), BG_CARD, 0.03)
    add_text(slide, Inches(1.1), y + Inches(0.1), Inches(2.5), Inches(0.35),
             label, 16, color, True)
    add_text(slide, Inches(3.8), y + Inches(0.1), Inches(8.5), Inches(0.7),
             desc, 14, LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 22 — TECH STACK
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "Technology Stack", 36, WHITE, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_CYAN)

categories = [
    ("Backend", ["FastAPI (Python)", "SQLAlchemy + PostgreSQL", "Celery + Redis", "Async/await throughout"], ACCENT_CYAN),
    ("Frontend", ["React 18", "Vite build tool", "Lucide icons", "Custom glassmorphic CSS"], ACCENT_PURPLE),
    ("AI / LLM", ["Groq API", "Llama 3.3 70B Versatile", "JSON-structured outputs", "Temperature-tuned per task"], ACCENT_GREEN),
    ("Infrastructure", ["Docker Compose", "Microservice architecture", "Gmail & Calendar APIs", "OAuth 2.0 authentication"], RGBColor(0xF5, 0x9E, 0x0B)),
]
for i, (cat, items, color) in enumerate(categories):
    x = Inches(0.8 + i * 3.15)
    y = Inches(1.6)
    add_shape_bg(slide, x, y, Inches(2.85), Inches(4.5), BG_CARD, 0.04)
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.4),
             cat, 20, color, True, PP_ALIGN.CENTER)
    add_accent_line(slide, x + Inches(0.5), y + Inches(0.65), Inches(1.85), color)
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(3.2),
                    items, 14, LIGHT, color, Pt(12))

# ════════════════════════════════════════════════════════════════
# SLIDE 23 — THANK YOU
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_line(slide, Inches(4.5), Inches(3.1), Inches(4.3), ACCENT_CYAN)
add_text(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(1.2),
         "Thank You", 54, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.8),
         "NORA EOS \u2014 ECHO \u2022 HERA \u2022 ARGUS", 22, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(4.5), Inches(13.333), Inches(0.6),
         "Questions?", 20, ACCENT_CYAN, False, PP_ALIGN.CENTER)


# ── Save ──
output_path = r"c:\Users\Shaggy\echo\NORA_EOS_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
